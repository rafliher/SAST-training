#!/usr/bin/env python3
"""
Add second trainer intro slide and LLM-as-SAST module to the presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import copy
import os

# ── Xecura Colors ──
XECURA_DARK_BLUE = RGBColor(0x1D, 0x5C, 0xA8)
XECURA_CYAN = RGBColor(0x00, 0xBA, 0xFF)
XECURA_RED = RGBColor(0xEF, 0x35, 0x4F)
XECURA_MEDIUM_BLUE = RGBColor(0x0F, 0x81, 0xC9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

SRC = 'day 1/Secure Coding - Xecura Branded.pptx'
DST = 'day 1/Secure Coding - Xecura Branded.pptx'


def add_run(para, text, font_name='Arial', font_size=Pt(12), bold=False, color=None, italic=False):
    """Add a formatted text run to a paragraph."""
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    if bold:
        run.font.bold = True
    if italic:
        run.font.italic = True
    if color:
        run.font.color.rgb = color
    return run


def add_textbox(slide, left, top, width, height, word_wrap=True):
    """Add a text box and return it."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = word_wrap
    return txBox


def clear_textframe(tf):
    """Clear all paragraphs from a text frame."""
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    tf.paragraphs[0].clear()


def clone_slide(prs, template_idx):
    """Clone a slide by index. Returns the new slide (appended at end)."""
    template = prs.slides[template_idx]
    slide_layout = template.slide_layout

    # Get the slide XML
    template_xml = copy.deepcopy(template._element)

    # Add a new slide with the same layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Clear the new slide
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy all elements from template
    for child in template_xml:
        new_slide._element.append(copy.deepcopy(child))

    # Copy relationships (images etc)
    for rel in template.part.rels.values():
        if "image" in rel.reltype:
            new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

    return new_slide


def move_slide(prs, old_index, new_index):
    """Move a slide from old_index to new_index (0-based)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    el = slides[old_index]
    xml_slides.remove(el)

    if new_index >= len(list(xml_slides)):
        xml_slides.append(el)
    else:
        ref = list(xml_slides)[new_index]
        xml_slides.insert(xml_slides.index(ref), el)


def create_trainer_slide(prs, insert_after_idx):
    """Create the second trainer INTROS slide for Herlambang Rafli Wicaksono."""
    print("Creating trainer intro slide for Rafli...")

    layout = prs.slide_masters[0].slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout)

    # Clear default placeholders
    for shape in list(slide.placeholders):
        sp = shape._element
        sp.getparent().remove(sp)

    # ── "INTROS" header (top-right) ──
    txb = add_textbox(slide, Emu(9090000), Emu(333846), Emu(3102000), Emu(800400))
    p = txb.text_frame.paragraphs[0]
    add_run(p, "INTROS", font_name='Arial', font_size=Pt(47), bold=True)

    # ── "Perkenalan Profile Trainer" subtitle ──
    txb = add_textbox(slide, Emu(7567912), Emu(989014), Emu(5307429), Emu(290464))
    p = txb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "Perkenalan Profile Trainer", font_size=Pt(18), color=XECURA_CYAN)

    # ── Name (bottom-left, over photo area) ──
    txb = add_textbox(slide, Emu(765078), Emu(5056079), Emu(3650806), Emu(800400))
    p = txb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "Herlambang Rafli W.", font_name='Arial', font_size=Pt(36), bold=True, color=XECURA_RED)

    # ── Left column: Education, Certifications, Experience ──
    txb = add_textbox(slide, Emu(4702286), Emu(115978), Emu(3925851), Emu(5700000))
    tf = txb.text_frame
    tf.word_wrap = True

    # Education
    p = tf.paragraphs[0]
    add_run(p, "Pendidikan :", font_size=Pt(12), bold=True)

    p = tf.add_paragraph()
    add_run(p, "Politeknik Siber dan Sandi Negara - Rekayasa Kriptografi (S.Tr.Kom)", font_size=Pt(12))

    # Blank line
    tf.add_paragraph()

    # Certifications
    p = tf.add_paragraph()
    add_run(p, "Sertifikasi :", font_size=Pt(12), bold=True)

    certs = [
        "eWPTX - INE Security",
        "CompTIA PenTest+ & Security+",
        "CompTIA CNVP",
        "TryHackMe PT1 Certified",
        "Junior Penetration Tester (BNSP)",
        "Certified AppSec Pentester (CAPen)",
        "Microsoft Security Operations Analyst",
        "Certified AppSec Practitioner (CAP)",
    ]
    for cert in certs:
        p = tf.add_paragraph()
        add_run(p, cert, font_size=Pt(12))

    # Blank line
    tf.add_paragraph()

    # Experience
    p = tf.add_paragraph()
    add_run(p, "Pengalaman:", font_size=Pt(12), bold=True)

    p = tf.add_paragraph()
    add_run(p, "Cybersecurity Researcher - Confidential Cyber Agency (2024-Now)", font_size=Pt(12), bold=True)

    experiences = [
        "R&D cybersecurity products, ISO 15408 compliance assessment",
        "Penetration testing for vital information infrastructure",
    ]
    for exp in experiences:
        p = tf.add_paragraph()
        add_run(p, exp, font_size=Pt(12))

    p = tf.add_paragraph()
    add_run(p, "Penetration Tester - PT Gan Mitra Usaha (2025-2026)", font_size=Pt(12), bold=True)

    p = tf.add_paragraph()
    add_run(p, "10,000+ vulns identified, 1,000+ assets managed", font_size=Pt(12))

    p = tf.add_paragraph()
    add_run(p, "Co-Founder - SiberLab ID (2023-Now)", font_size=Pt(12), bold=True)

    p = tf.add_paragraph()
    add_run(p, "10+ products, 10+ national CTF events, 500+ students", font_size=Pt(12))

    # ── Right column: Achievements & Publications ──
    txb = add_textbox(slide, Emu(8665358), Emu(1558496), Emu(3415749), Emu(4500781))
    tf = txb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    add_run(p, "Penghargaan :", font_size=Pt(12), bold=True)

    achievements = [
        "Finalis Cyber Jawara 2023",
        "6th Place ASCIS CTF (ASEAN)",
        "3rd Place Networking TECHNOLOGY EUPHORIA 2023",
        "Finalist HOLOGY 6.0 CTF",
        "Finalist Techcompfest CTF 2023",
        "Finalist NU-TECH CTF 2022",
    ]
    for ach in achievements:
        p = tf.add_paragraph()
        add_run(p, ach, font_size=Pt(12))

    tf.add_paragraph()
    p = tf.add_paragraph()
    add_run(p, "Publikasi :", font_size=Pt(12), bold=True)

    publications = [
        "SiberBox: Lightweight Sandbox (EECSI 2025)",
        "SSO Implementation (JEPIN 2025)",
        "Spaticrypt: Crypto Education Platform (2025)",
        "Secure SDLC Framework (ICICoS, IEEE 2024)",
        "Sentiment Analysis with GPT-4 (ICIMCIS 2024)",
        "Phishing URL Detection Benchmark (ICIMCIS 2023)",
    ]
    for pub in publications:
        p = tf.add_paragraph()
        add_run(p, pub, font_size=Pt(12))

    # Move to correct position (after slide 5)
    total = len(prs.slides)
    move_slide(prs, total - 1, insert_after_idx)
    print(f"  Inserted trainer slide at position {insert_after_idx + 1}")


def create_module_divider(prs, module_num, title, insert_idx):
    """Create a module divider slide matching the existing style."""
    # Use the Title Only layout
    layout = prs.slide_masters[0].slide_layouts[3]  # Title Only
    slide = prs.slides.add_slide(layout)

    # Clear default placeholders content
    for shape in list(slide.placeholders):
        if shape.placeholder_format.idx == 0:  # Title
            sp = shape._element
            sp.getparent().remove(sp)

    # Add the module divider background image (reuse from existing slide 9)
    # We'll copy the image from the original slide 9
    src_slide = prs.slides[9]  # slide 10 now (after trainer insert), was slide 9
    for shape in src_slide.shapes:
        if shape.shape_type == 13 and shape.name == 'Picture 2':
            # Get the image blob
            img_blob = shape.image.blob
            img_ct = shape.image.content_type
            # Save temp
            ext = 'png' if 'png' in img_ct else 'jpg'
            tmp = f'/tmp/module_bg.{ext}'
            with open(tmp, 'wb') as f:
                f.write(img_blob)
            # Add to new slide
            slide.shapes.add_picture(
                tmp, Emu(756541), Emu(0), Emu(5129561), Emu(6858000)
            )
            break

    # Add the rounded rectangle with module title
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUND_2_DIAG_RECTANGLE,
        Emu(6665794), Emu(1572322),
        Emu(4432725), Emu(3713356)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = XECURA_DARK_BLUE
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    add_run(p, f"Module {module_num}:", font_name='Arial', font_size=Pt(40), color=WHITE)

    p = tf.add_paragraph()
    add_run(p, title, font_name='Arial', font_size=Pt(40), color=WHITE)

    # Move to position
    total = len(prs.slides)
    move_slide(prs, total - 1, insert_idx)
    print(f"  Created module divider: Module {module_num} at position {insert_idx + 1}")
    return insert_idx


def create_content_slide(prs, title_text, body_content, insert_idx, title_size=Pt(40)):
    """Create a content slide with title and body text.

    body_content is a list of tuples: (text, bold, font_size, color, indent_level)
    """
    layout = prs.slide_masters[0].slide_layouts[3]  # Title Only
    slide = prs.slides.add_slide(layout)

    # Set title
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:  # Title
            shape.text = title_text
            for run in shape.text_frame.paragraphs[0].runs:
                run.font.size = title_size
                run.font.bold = True
            break

    # Add body text box
    txb = add_textbox(slide,
                      Emu(700000), Emu(1200000),
                      Emu(10800000), Emu(5200000))
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for item in body_content:
        if isinstance(item, str):
            # Simple text
            text, bold, size, color, indent = item, False, Pt(22), None, 0
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            size = item[2] if len(item) > 2 else Pt(22)
            color = item[3] if len(item) > 3 else None
            indent = item[4] if len(item) > 4 else 0

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        if indent > 0:
            p.level = indent

        add_run(p, text, font_size=size, bold=bold, color=color)

    # Move to position
    total = len(prs.slides)
    move_slide(prs, total - 1, insert_idx)
    return slide


def create_code_slide(prs, title_text, code_blocks, insert_idx):
    """Create a slide with code examples.

    code_blocks is a list of (label, code_text)
    """
    layout = prs.slide_masters[0].slide_layouts[3]
    slide = prs.slides.add_slide(layout)

    # Title
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title_text
            for run in shape.text_frame.paragraphs[0].runs:
                run.font.size = Pt(36)
                run.font.bold = True
            break

    # Code blocks
    y_pos = Emu(1200000)
    for label, code_text in code_blocks:
        # Label
        txb = add_textbox(slide, Emu(700000), y_pos, Emu(10800000), Emu(300000))
        p = txb.text_frame.paragraphs[0]
        add_run(p, label, font_size=Pt(16), bold=True, color=XECURA_DARK_BLUE)
        y_pos += Emu(350000)

        # Code box
        code_height = Emu(min(len(code_text.split('\n')) * 200000 + 150000, 2500000))
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(700000), y_pos,
            Emu(10800000), code_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)  # Dark code bg
        shape.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(150000)
        tf.margin_top = Emu(100000)

        for i, line in enumerate(code_text.strip().split('\n')):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            add_run(p, line, font_name='Consolas', font_size=Pt(11),
                    color=RGBColor(0xCC, 0xDD, 0xEE))

        y_pos += code_height + Emu(200000)

    total = len(prs.slides)
    move_slide(prs, total - 1, insert_idx)
    return slide


def create_llm_sast_module(prs, insert_idx):
    """Create the complete LLM-as-SAST module."""
    print("Creating LLM-as-SAST module slides...")
    idx = insert_idx

    # ── Slide 1: Module Divider ──
    create_module_divider(prs, "7", "LLM-Powered Source Code\nAnalysis (AI-SAST)", idx)
    idx += 1

    # ── Slide 2: Why LLM for SAST? ──
    create_content_slide(prs, "Why LLM for SAST?", [
        ("Traditional SAST Limitations:", True, Pt(22), XECURA_DARK_BLUE),
        ("Rule-based pattern matching → high false positive rates", False, Pt(20), None, 1),
        ("Cannot understand business logic or contextual vulnerabilities", False, Pt(20), None, 1),
        ("Limited to known vulnerability patterns (signatures)", False, Pt(20), None, 1),
        ("Struggles with framework-specific code and custom abstractions", False, Pt(20), None, 1),
        ("", False, Pt(12)),
        ("LLM-Powered SAST Advantages:", True, Pt(22), XECURA_DARK_BLUE),
        ("Semantic understanding of code logic and data flow", False, Pt(20), None, 1),
        ("Can reason about business logic flaws and design issues", False, Pt(20), None, 1),
        ("Understands context: framework conventions, API patterns, auth flows", False, Pt(20), None, 1),
        ("Provides human-readable explanations and fix suggestions", False, Pt(20), None, 1),
        ("Adapts to new languages/frameworks without new rules", False, Pt(20), None, 1),
    ], idx)
    idx += 1

    # ── Slide 3: Traditional vs LLM SAST Comparison ──
    create_content_slide(prs, "Traditional SAST vs LLM-Powered SAST", [
        ("Aspect               | Traditional SAST      | LLM-Powered SAST", True, Pt(18), XECURA_DARK_BLUE),
        ("─────────────────────────────────────────────────────────────────", False, Pt(14)),
        ("Detection Method  | Regex / AST patterns      | Semantic code understanding", False, Pt(17)),
        ("False Positives      | High (30-70%)                | Low (contextual filtering)", False, Pt(17)),
        ("Business Logic      | Cannot detect                | Can reason about logic flaws", False, Pt(17)),
        ("Explanation            | Generic CWE reference    | Natural language explanation", False, Pt(17)),
        ("Fix Suggestions     | Template-based              | Context-aware code fixes", False, Pt(17)),
        ("New Frameworks   | Requires new rules          | Learns from training data", False, Pt(17)),
        ("Speed                    | Very fast (ms)                | Slower (seconds per file)", False, Pt(17)),
        ("Offline Support     | Full offline                       | Depends on model hosting", False, Pt(17)),
        ("", False, Pt(12)),
        ("Best Practice: Use BOTH together — traditional SAST for speed + LLM for depth", True, Pt(18), XECURA_RED),
    ], idx, title_size=Pt(36))
    idx += 1

    # ── Slide 4: How LLM Analyzes Code ──
    create_content_slide(prs, "How LLM Analyzes Source Code", [
        ("1. Code Chunking & Context Window", True, Pt(22), XECURA_DARK_BLUE),
        ("   Source code is split into manageable chunks (functions, classes, files)", False, Pt(18)),
        ("   Context window includes imports, dependencies, and related code", False, Pt(18)),
        ("", False, Pt(10)),
        ("2. Prompt Engineering for Security Analysis", True, Pt(22), XECURA_DARK_BLUE),
        ("   System prompt defines the security expert persona and analysis rules", False, Pt(18)),
        ("   Structured output format (JSON) for consistent vulnerability reports", False, Pt(18)),
        ("", False, Pt(10)),
        ("3. Multi-Pass Analysis", True, Pt(22), XECURA_DARK_BLUE),
        ("   Pass 1: Identify potential vulnerability locations", False, Pt(18)),
        ("   Pass 2: Validate with data flow and taint analysis reasoning", False, Pt(18)),
        ("   Pass 3: Generate fix recommendations with secure code examples", False, Pt(18)),
        ("", False, Pt(10)),
        ("4. Integration with CI/CD Pipeline", True, Pt(22), XECURA_DARK_BLUE),
        ("   Triggered on pull requests, commit hooks, or scheduled scans", False, Pt(18)),
    ], idx)
    idx += 1

    # ── Slide 5: Available LLM Tools for SAST ──
    create_content_slide(prs, "LLM-Powered SAST Tools & Models", [
        ("Commercial / SaaS:", True, Pt(22), XECURA_DARK_BLUE),
        ("GitHub Copilot (Autofix) — integrated in GitHub Advanced Security", False, Pt(18), None, 1),
        ("Snyk DeepCode AI — AI-powered code analysis in Snyk platform", False, Pt(18), None, 1),
        ("Amazon CodeGuru — ML-based code review for Java/Python", False, Pt(18), None, 1),
        ("Semgrep + AI Assistant — rule-based + LLM-assisted analysis", False, Pt(18), None, 1),
        ("", False, Pt(10)),
        ("Open Source / Self-Hosted:", True, Pt(22), XECURA_DARK_BLUE),
        ("Claude Code / Claude API — Anthropic's model, strong at code analysis", False, Pt(18), None, 1),
        ("GPT-4 / OpenAI API — general purpose, effective with good prompts", False, Pt(18), None, 1),
        ("Ollama + CodeLlama/DeepSeek — fully offline, self-hosted LLM", False, Pt(18), None, 1),
        ("Aider / Continue.dev — IDE-integrated AI code review", False, Pt(18), None, 1),
        ("", False, Pt(10)),
        ("For this training: We use Ollama (offline) + Claude API (cloud)", True, Pt(20), XECURA_RED),
    ], idx)
    idx += 1

    # ── Slide 6: Setup Ollama ──
    create_content_slide(prs, "Setup: Ollama for Offline LLM-SAST", [
        ("Step 1: Install Ollama", True, Pt(22), XECURA_DARK_BLUE),
        ("   curl -fsSL https://ollama.com/install.sh | sh", False, Pt(16)),
        ("   # Or download from https://ollama.com for macOS/Windows", False, Pt(16)),
        ("", False, Pt(10)),
        ("Step 2: Pull a Code Analysis Model", True, Pt(22), XECURA_DARK_BLUE),
        ("   ollama pull deepseek-coder-v2:16b    # Best for code analysis", False, Pt(16)),
        ("   ollama pull codellama:13b             # Alternative: Meta's CodeLlama", False, Pt(16)),
        ("   ollama pull qwen2.5-coder:14b         # Alternative: Alibaba's Qwen", False, Pt(16)),
        ("", False, Pt(10)),
        ("Step 3: Verify Installation", True, Pt(22), XECURA_DARK_BLUE),
        ("   ollama list                            # Check downloaded models", False, Pt(16)),
        ("   ollama run deepseek-coder-v2:16b      # Test interactive chat", False, Pt(16)),
        ("", False, Pt(10)),
        ("Requirements: 16GB+ RAM for 13B-16B models, GPU recommended", True, Pt(18), XECURA_RED),
        ("For lightweight: ollama pull deepseek-coder-v2:lite (needs ~4GB RAM)", False, Pt(16), XECURA_MEDIUM_BLUE),
    ], idx)
    idx += 1

    # ── Slide 7: Prompt Engineering for Security ──
    create_code_slide(prs, "Prompt Engineering: Security Analysis Prompt", [
        ("System Prompt Template:", """You are an expert application security engineer.
Analyze the following source code for security vulnerabilities.

For each vulnerability found, provide:
1. Vulnerability type (e.g., SQL Injection, XSS, IDOR)
2. CWE ID
3. Severity (Critical/High/Medium/Low)
4. Affected line numbers
5. Explanation of the attack vector
6. Secure code fix

Focus on OWASP Top 10 categories. Only report real
vulnerabilities, not style issues. Be precise."""),
        ("User Prompt Template:", """Analyze this PHP code for security vulnerabilities:

```php
// [paste target source code here]
```

Respond in JSON format with a "vulnerabilities" array."""),
    ], idx)
    idx += 1

    # ── Slide 8: Python Script for LLM-SAST ──
    create_code_slide(prs, "Building an LLM-SAST Scanner (Python + Ollama)", [
        ("llm_sast_scanner.py:", """import ollama, json, sys, pathlib

SYSTEM_PROMPT = pathlib.Path("prompt.txt").read_text()

def scan_file(filepath):
    code = pathlib.Path(filepath).read_text()
    response = ollama.chat(model='deepseek-coder-v2:16b',
        messages=[
            {'role': 'system', 'content': SYSTEM_PROMPT},
            {'role': 'user',
             'content': f'Analyze this code:\\n```\\n{code}\\n```'}
        ], format='json')
    return json.loads(response['message']['content'])

if __name__ == '__main__':
    results = scan_file(sys.argv[1])
    for vuln in results.get('vulnerabilities', []):
        sev = vuln.get('severity','?')
        typ = vuln.get('type','?')
        line = vuln.get('line','?')
        print(f'[{sev}] {typ} at line {line}')
        print(f'  → {vuln.get("explanation","")}')"""),
    ], idx)
    idx += 1

    # ── Slide 9: Usage & Running the scanner ──
    create_code_slide(prs, "Running the LLM-SAST Scanner", [
        ("Basic Usage:", """# Scan a single file
python llm_sast_scanner.py vulnerable_app.php

# Scan entire project directory
find ./src -name "*.php" -exec python llm_sast_scanner.py {} \\;

# Output example:
[CRITICAL] SQL Injection at line 23
  → User input $_GET['id'] directly concatenated into SQL query
    without parameterization. Use PDO prepared statements.
[HIGH] XSS (Stored) at line 45
  → User-supplied data echoed without htmlspecialchars().
    Apply output encoding before rendering in HTML context.
[MEDIUM] Insecure Direct Object Reference at line 67
  → User can modify 'user_id' parameter to access other
    users' data. Implement ownership validation."""),
    ], idx)
    idx += 1

    # ── Slide 10: Claude API for SAST ──
    create_code_slide(prs, "Using Claude API for Cloud-Based LLM-SAST", [
        ("claude_sast.py:", """import anthropic, json, sys, pathlib

client = anthropic.Anthropic()  # uses ANTHROPIC_API_KEY env var
SYSTEM_PROMPT = pathlib.Path("prompt.txt").read_text()

def scan_file(filepath):
    code = pathlib.Path(filepath).read_text()
    message = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4096,
        system=SYSTEM_PROMPT,
        messages=[{'role': 'user',
            'content': f'Analyze:\\n```\\n{code}\\n```'}])
    return json.loads(message.content[0].text)

# Setup: pip install anthropic
# export ANTHROPIC_API_KEY="sk-ant-..."
# python claude_sast.py target_app.php"""),
    ], idx)
    idx += 1

    # ── Slide 11: CI/CD Integration ──
    create_code_slide(prs, "CI/CD Integration: GitHub Actions", [
        (".github/workflows/llm-sast.yml:", """name: LLM Security Scan
on: [pull_request]
jobs:
  llm-sast:
    runs-on: ubuntu-latest
    steps:
      - uses: actions/checkout@v4
      - uses: actions/setup-python@v5
        with: { python-version: '3.12' }
      - run: pip install anthropic
      - name: Run LLM-SAST on changed files
        env:
          ANTHROPIC_API_KEY: ${{ secrets.ANTHROPIC_API_KEY }}
        run: |
          git diff --name-only origin/main...HEAD \\
            --diff-filter=ACMR -- '*.php' '*.py' '*.js' |
          while read file; do
            echo "Scanning: $file"
            python llm_sast_scanner.py "$file" >> report.md
          done
      - name: Post results as PR comment
        uses: actions/github-script@v7
        with:
          script: |
            const fs = require('fs');
            const body = fs.readFileSync('report.md','utf8');
            github.rest.issues.createComment({
              ...context.repo, issue_number: context.issue.number,
              body: '## LLM-SAST Results\\n' + body });"""),
    ], idx)
    idx += 1

    # ── Slide 12: Live Demo Scenario ──
    create_content_slide(prs, "Hands-On: LLM vs Traditional SAST", [
        ("Demo Scenario: Scan DVWA (Damn Vulnerable Web Application)", True, Pt(22), XECURA_DARK_BLUE),
        ("", False, Pt(10)),
        ("Step 1: Run SonarQube scan on DVWA source code", False, Pt(20)),
        ("   → Observe: generic findings, many false positives", False, Pt(18), XECURA_MEDIUM_BLUE, 1),
        ("", False, Pt(10)),
        ("Step 2: Run our LLM-SAST scanner on the same files", False, Pt(20)),
        ("   → Observe: contextual findings with explanations", False, Pt(18), XECURA_MEDIUM_BLUE, 1),
        ("", False, Pt(10)),
        ("Step 3: Compare results side-by-side", False, Pt(20)),
        ("   → Which found real vulns? Which had fewer false positives?", False, Pt(18), XECURA_MEDIUM_BLUE, 1),
        ("", False, Pt(10)),
        ("Step 4: Use LLM to generate secure patches", False, Pt(20)),
        ("   → Ask the LLM to rewrite vulnerable functions securely", False, Pt(18), XECURA_MEDIUM_BLUE, 1),
        ("", False, Pt(10)),
        ("Step 5: Re-scan patched code to verify fixes", False, Pt(20)),
        ("   → Confirm the vulnerability is resolved", False, Pt(18), XECURA_MEDIUM_BLUE, 1),
    ], idx)
    idx += 1

    # ── Slide 13: Best Practices ──
    create_content_slide(prs, "LLM-SAST Best Practices & Limitations", [
        ("Best Practices:", True, Pt(24), XECURA_DARK_BLUE),
        ("Combine with traditional SAST — use LLM as a second opinion layer", False, Pt(20), None, 1),
        ("Use structured output (JSON) for automated processing", False, Pt(20), None, 1),
        ("Version your prompts — treat them as code (prompt-as-code)", False, Pt(20), None, 1),
        ("Review LLM findings manually — AI can hallucinate vulnerabilities", False, Pt(20), None, 1),
        ("Use self-hosted models (Ollama) for sensitive/proprietary code", False, Pt(20), None, 1),
        ("Keep models updated — newer models = better security understanding", False, Pt(20), None, 1),
        ("", False, Pt(10)),
        ("Limitations to Be Aware Of:", True, Pt(24), XECURA_RED),
        ("Context window limits — large codebases need chunking strategies", False, Pt(20), None, 1),
        ("Hallucination risk — LLM may report non-existent vulnerabilities", False, Pt(20), None, 1),
        ("Not deterministic — same code may produce different results", False, Pt(20), None, 1),
        ("Speed — slower than rule-based scanners for large projects", False, Pt(20), None, 1),
        ("Cost — cloud API usage incurs per-token charges", False, Pt(20), None, 1),
    ], idx)
    idx += 1

    # ── Slide 14: Summary ──
    create_content_slide(prs, "Module Summary: LLM-Powered SAST", [
        ("Key Takeaways:", True, Pt(24), XECURA_DARK_BLUE),
        ("", False, Pt(10)),
        ("1.  LLMs bring semantic understanding to source code analysis", False, Pt(22)),
        ("     — beyond pattern matching to real comprehension", False, Pt(18), XECURA_MEDIUM_BLUE),
        ("", False, Pt(10)),
        ("2.  Best used as a complement to traditional SAST, not a replacement", False, Pt(22)),
        ("     — defense in depth applies to tooling too", False, Pt(18), XECURA_MEDIUM_BLUE),
        ("", False, Pt(10)),
        ("3.  Ollama enables fully offline, private code analysis", False, Pt(22)),
        ("     — no code leaves your infrastructure", False, Pt(18), XECURA_MEDIUM_BLUE),
        ("", False, Pt(10)),
        ("4.  Prompt engineering is critical for accurate security analysis", False, Pt(22)),
        ("     — structured prompts = structured, actionable results", False, Pt(18), XECURA_MEDIUM_BLUE),
        ("", False, Pt(10)),
        ("5.  CI/CD integration makes LLM-SAST part of your SDLC", False, Pt(22)),
        ("     — shift-left with AI-powered security review", False, Pt(18), XECURA_MEDIUM_BLUE),
    ], idx)
    idx += 1

    return idx  # Return next available index


def update_outline_slide(prs):
    """Update the outline slide (slide 8) to include the new LLM module."""
    print("Updating outline slide...")

    # Find the outline slide (has "Outline" title)
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text_frame.paragraphs[0].text.strip() == 'Outline':
                    # Found it - now find the body textbox
                    for s in slide.shapes:
                        if s.shape_type == 17 and s.name == 'object 2':  # TextBox
                            tf = s.text_frame
                            # Check if already has LLM module
                            for p in tf.paragraphs:
                                if 'LLM' in p.text:
                                    print("  Outline already has LLM module")
                                    return

                            # Add new module to outline
                            p = tf.add_paragraph()
                            add_run(p, "LLM-Powered Source Code Analysis (AI-SAST)",
                                    font_size=Pt(24), color=XECURA_RED, bold=True)
                            print(f"  Added LLM module to outline (slide {i+1})")
                            return


def renumber_modules(prs):
    """Renumber existing Module 7, 8, 9 to 8, 9, 10 since we inserted Module 7."""
    print("Renumbering existing modules...")
    renumber_map = {
        "Module 7:": "Module 8:",
        "Module 8:": "Module 9:",
        "Module 9:": "Module 10:",
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for old, new in renumber_map.items():
                            if run.text.strip().startswith(old.rstrip(':')):
                                run.text = run.text.replace(old.rstrip(':'), new.rstrip(':'))
                                print(f"  Renumbered: {old} → {new}")


def main():
    print(f"Loading {SRC}...")
    prs = Presentation(SRC)
    total_before = len(prs.slides)
    print(f"  {total_before} slides loaded")

    # Step 1: Add trainer slide after slide 5 (index 5)
    create_trainer_slide(prs, insert_after_idx=5)

    # Step 2: Renumber existing modules 7→8, 8→9
    # The existing Module 7 (Hands-On Lab) is now at a shifted position
    # After inserting trainer slide, everything shifts by 1
    # Module 6 (Security Toolchain) was at slide 83, now at 84
    # We want to insert our new Module 7 AFTER Module 6's last slide

    # Find where Module 6 ends (look for the next module divider or specific slide)
    # Module 6 slides: 84-92 (after trainer insert shift)
    # Let's find it dynamically
    module7_insert_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'Module 7:' in run.text or (run.text.strip().startswith('Module 7') and ':' in para.text):
                            module7_insert_idx = i
                            break

    if module7_insert_idx is None:
        # Fallback: insert before "Hands-On Lab" section
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if 'Hands-On Lab' in shape.text_frame.text or 'Hands-on Lab' in shape.text_frame.text:
                        module7_insert_idx = i
                        break
            if module7_insert_idx:
                break

    if module7_insert_idx is None:
        module7_insert_idx = 92  # Fallback position

    print(f"\nInserting LLM module at position {module7_insert_idx + 1}")

    # Step 3: Renumber existing modules BEFORE inserting
    renumber_modules(prs)

    # Step 4: Create LLM-as-SAST module
    create_llm_sast_module(prs, module7_insert_idx)

    # Step 5: Update outline slide
    update_outline_slide(prs)

    total_after = len(prs.slides)
    print(f"\nTotal slides: {total_before} → {total_after} (+{total_after - total_before})")
    print(f"Saving to {DST}...")
    prs.save(DST)
    print("Done!")


if __name__ == '__main__':
    main()
