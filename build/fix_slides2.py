#!/usr/bin/env python3
"""
Fix round 2:
1. Fix all generated slide title positions (move to top, full width)
2. Replace OWASP 2021 comparison image with 2025 version
3. Fix text body positions accordingly
"""

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
import os

XECURA_DARK_BLUE = RGBColor(0x1D, 0x5C, 0xA8)
XECURA_CYAN = RGBColor(0x00, 0xBA, 0xFF)
XECURA_RED = RGBColor(0xEF, 0x35, 0x4F)
XECURA_MEDIUM_BLUE = RGBColor(0x0F, 0x81, 0xC9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SRC = 'day 1/Secure Coding - Xecura Branded.pptx'
DST = 'day 1/Secure Coding - Xecura Branded.pptx'

# Target title position (matching original slides like slide 11, 41)
TITLE_LEFT = Emu(813568)
TITLE_TOP = Emu(320278)
TITLE_WIDTH = Emu(11296656)
TITLE_HEIGHT = Emu(755982)


def add_run(para, text, font_name='Arial', font_size=Pt(12), bold=False, color=None):
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color
    return run


def generate_owasp_2025_image():
    """Generate a clean OWASP Top 10 2025 comparison chart image."""
    w, h = 1800, 900
    img = Image.new('RGB', (w, h), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Try to use a decent font
    try:
        font_title = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 28)
        font_header = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
        font_body = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
        font_small = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 13)
    except:
        font_title = ImageFont.load_default()
        font_header = font_title
        font_body = font_title
        font_small = font_title

    # Colors
    blue = (29, 92, 168)
    cyan = (0, 186, 255)
    red = (239, 53, 79)
    dark = (51, 51, 51)
    light_gray = (240, 240, 245)
    white = (255, 255, 255)
    green_new = (0, 168, 120)

    # Title
    draw.text((w // 2 - 200, 15), "OWASP Top 10 — 2021 vs 2025", fill=blue, font=font_title)

    # Column headers
    col1_x = 50
    col2_x = 950
    header_y = 65

    # Left column header
    draw.rounded_rectangle([col1_x - 10, header_y, col1_x + 850, header_y + 35], radius=5, fill=blue)
    draw.text((col1_x + 300, header_y + 5), "2021", fill=white, font=font_header)

    # Right column header
    draw.rounded_rectangle([col2_x - 10, header_y, col2_x + 850, header_y + 35], radius=5, fill=red)
    draw.text((col2_x + 300, header_y + 5), "2025", fill=white, font=font_header)

    items_2021 = [
        "A01  Broken Access Control",
        "A02  Cryptographic Failures",
        "A03  Injection",
        "A04  Insecure Design",
        "A05  Security Misconfiguration",
        "A06  Vulnerable & Outdated Components",
        "A07  Identification & Auth Failures",
        "A08  Software & Data Integrity Failures",
        "A09  Security Logging & Monitoring",
        "A10  Server-Side Request Forgery",
    ]
    items_2025 = [
        "A01  Broken Access Control",
        "A02  Security Misconfiguration",
        "A03  Software Supply Chain Failures",
        "A04  Cryptographic Failures",
        "A05  Injection",
        "A06  Insecure Design",
        "A07  Authentication Failures",
        "A08  Software or Data Integrity Failures",
        "A09  Security Logging & Alerting Failures",
        "A10  Mishandling of Exceptional Conditions",
    ]

    # New entries in 2025
    new_in_2025 = {"A03", "A10"}

    y_start = 115
    row_h = 72
    row_padding = 5

    for i in range(10):
        y = y_start + i * row_h

        # Left (2021)
        bg_l = light_gray if i % 2 == 0 else white
        draw.rounded_rectangle([col1_x - 10, y, col1_x + 850, y + row_h - row_padding], radius=8, fill=bg_l)

        # Code badge
        code_2021 = items_2021[i][:3]
        draw.rounded_rectangle([col1_x, y + 10, col1_x + 45, y + row_h - 15], radius=4, fill=blue)
        draw.text((col1_x + 5, y + 16), code_2021, fill=white, font=font_small)
        draw.text((col1_x + 55, y + 18), items_2021[i][5:], fill=dark, font=font_body)

        # Right (2025)
        code_2025 = items_2025[i][:3]
        is_new = code_2025 in new_in_2025
        bg_r = (255, 240, 240) if is_new else (light_gray if i % 2 == 0 else white)
        draw.rounded_rectangle([col2_x - 10, y, col2_x + 850, y + row_h - row_padding], radius=8, fill=bg_r)

        badge_color = green_new if is_new else red
        draw.rounded_rectangle([col2_x, y + 10, col2_x + 45, y + row_h - 15], radius=4, fill=badge_color)
        draw.text((col2_x + 5, y + 16), code_2025, fill=white, font=font_small)
        draw.text((col2_x + 55, y + 18), items_2025[i][5:], fill=dark, font=font_body)

        if is_new:
            draw.text((col2_x + 750, y + 18), "NEW", fill=green_new, font=font_small)

    # Arrow hints for major moves
    # A05:2021 (Security Misconfig) → A02:2025
    draw.line([(col1_x + 855, y_start + 4 * row_h + row_h // 2),
               (col2_x - 15, y_start + 1 * row_h + row_h // 2)],
              fill=cyan, width=2)
    # A02:2021 (Crypto) → A04:2025
    draw.line([(col1_x + 855, y_start + 1 * row_h + row_h // 2),
               (col2_x - 15, y_start + 3 * row_h + row_h // 2)],
              fill=cyan, width=2)

    # Legend
    ly = h - 45
    draw.rounded_rectangle([col1_x, ly, col1_x + 20, ly + 20], radius=3, fill=green_new)
    draw.text((col1_x + 25, ly), "= New in 2025", fill=dark, font=font_small)
    draw.text((col1_x + 200, ly), "Lines show major position changes", fill=(150, 150, 150), font=font_small)

    path = '/tmp/owasp_2025_chart.png'
    img.save(path, 'PNG')
    return path


def main():
    print(f"Loading {SRC}...")
    prs = Presentation(SRC)

    # ═══════════════════════════════════
    # FIX 1: Title positions on all generated slides
    # ═══════════════════════════════════
    print("\n--- Fix 1: Title positions ---")

    # Find all our generated slides (94-105) and fix their Title placeholders
    for si in range(93, min(106, len(prs.slides))):
        slide = prs.slides[si]
        for shape in slide.shapes:
            # Fix Title placeholders that are in the wrong position
            if shape.shape_type == 14 and shape.name == 'Title 1':
                old_pos = (shape.left, shape.top, shape.width, shape.height)
                # Only fix if it's in the bad center position
                if shape.top > Emu(1000000):
                    shape.left = TITLE_LEFT
                    shape.top = TITLE_TOP
                    shape.width = TITLE_WIDTH
                    shape.height = TITLE_HEIGHT

                    # Also fix the text alignment
                    for para in shape.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.LEFT

                    title_text = shape.text_frame.text[:50]
                    print(f"  Slide {si+1}: Fixed title \"{title_text}\"")

    # ═══════════════════════════════════
    # FIX 2: Replace OWASP comparison image (slide 40)
    # ═══════════════════════════════════
    print("\n--- Fix 2: OWASP 2025 image ---")
    chart_path = generate_owasp_2025_image()
    print(f"  Generated chart: {chart_path}")

    slide40 = prs.slides[39]
    for shape in slide40.shapes:
        if shape.shape_type == 13:  # Picture
            # Replace the image
            old_left, old_top = shape.left, shape.top
            old_w, old_h = shape.width, shape.height

            # Get the blip element and replace via relationship
            pic = shape._element
            blip = pic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            if blip is not None:
                embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if embed:
                    target_part = slide40.part.rels[embed].target_part
                    with open(chart_path, 'rb') as f:
                        target_part._blob = f.read()
                    print(f"  Replaced OWASP image on slide 40")

            # Adjust size to better fit
            shape.left = Emu(813568)
            shape.top = Emu(1200000)
            shape.width = Emu(10500000)
            shape.height = Emu(5250000)
            break

    # Also update the text on slide 40
    for shape in slide40.shapes:
        if shape.has_text_frame:
            full = shape.text_frame.text
            if 'standard awareness' in full:
                # This is the description textbox - update it
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'most critical security risks' in run.text:
                            run.text = run.text  # keep as is
                        run.text = run.text.replace('2021', '2025').replace('2025)', '2025)')

    # ═══════════════════════════════════
    # FIX 3: Also fix body textbox positions on generated slides
    # The body textboxes may overlap with the repositioned titles
    # ═══════════════════════════════════
    print("\n--- Fix 3: Body text positions ---")
    # Our generated slides use TextBox shapes for body content
    # The title is now at top (320278), so body should start below it
    BODY_TOP = Emu(1100000)

    for si in range(93, min(106, len(prs.slides))):
        slide = prs.slides[si]
        for shape in slide.shapes:
            if shape.shape_type == 17:  # TextBox
                # Check if this is a body textbox (large, contains multiple paragraphs)
                lines = sum(1 for p in shape.text_frame.paragraphs if p.text.strip())
                if lines >= 3 and shape.height > Emu(2000000):
                    if shape.top < BODY_TOP:
                        shape.top = BODY_TOP
                        print(f"  Slide {si+1}: Moved body text down")

    # ═══════════════════════════════════
    # SAVE
    # ═══════════════════════════════════
    print(f"\nSaving to {DST}...")
    prs.save(DST)
    print("Done!")


if __name__ == '__main__':
    main()
