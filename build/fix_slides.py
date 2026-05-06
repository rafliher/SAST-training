#!/usr/bin/env python3
"""
Fix all slide issues:
1. Fix positioning/overflow on generated slides
2. Rearrange hands-on into Module 8
3. Update OWASP 2021 → 2025
4. Update Ollama model references to match what's installed
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy

XECURA_DARK_BLUE = RGBColor(0x1D, 0x5C, 0xA8)
XECURA_CYAN = RGBColor(0x00, 0xBA, 0xFF)
XECURA_RED = RGBColor(0xEF, 0x35, 0x4F)
XECURA_MEDIUM_BLUE = RGBColor(0x0F, 0x81, 0xC9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SRC = 'day 1/Secure Coding - Xecura Branded.pptx'
DST = 'day 1/Secure Coding - Xecura Branded.pptx'

SLIDE_W = 12192000
SLIDE_H = 6858000


def add_run(para, text, font_name='Arial', font_size=Pt(12), bold=False, color=None, italic=False):
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    if bold:
        run.font.bold = True
    if italic:
        run.font.italic = True
    if color:
        run.font.color.rgb = color
    return run


def add_textbox(slide, left, top, width, height, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = word_wrap
    return txBox


def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    el = slides[old_index]
    xml_slides.remove(el)
    if new_index >= len(list(xml_slides)):
        xml_slides.append(el)
    else:
        ref = list(xml_slides)[new_index]
        xml_slides.insert(xml_slides.index(ref), el)


def delete_slide(prs, index):
    """Delete slide at given index."""
    rId = prs.slides._sldIdLst[index].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    sldId = prs.slides._sldIdLst[index]
    prs.slides._sldIdLst.remove(sldId)
    # Also remove the relationship
    if rId:
        try:
            del prs.part.rels[rId]
        except KeyError:
            pass


def clear_slide_shapes(slide, keep_placeholders=False):
    """Remove all shapes from a slide."""
    for shape in list(slide.shapes):
        if keep_placeholders and shape.shape_type == 14:
            continue
        sp = shape._element
        sp.getparent().remove(sp)


def rebuild_code_slide(slide, title_text, code_blocks):
    """Rebuild a code slide with proper sizing."""
    clear_slide_shapes(slide, keep_placeholders=True)

    # Set title in placeholder if exists
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title_text
            for run in shape.text_frame.paragraphs[0].runs:
                run.font.size = Pt(32)
                run.font.bold = True
            break

    y_pos = Emu(1150000)
    for label, code_text in code_blocks:
        lines = code_text.strip().split('\n')

        # Label
        txb = add_textbox(slide, Emu(500000), y_pos, Emu(11200000), Emu(260000))
        p = txb.text_frame.paragraphs[0]
        add_run(p, label, font_size=Pt(14), bold=True, color=XECURA_DARK_BLUE)
        y_pos += Emu(280000)

        # Code box - calculate height properly
        line_h = Emu(175000)  # per line at Pt(10)
        padding = Emu(180000)
        code_height = min(len(lines) * line_h + padding, Emu(SLIDE_H - 1400000) - y_pos + Emu(1150000))
        code_height = max(code_height, Emu(500000))

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(500000), y_pos,
            Emu(11200000), code_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
        shape.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(120000)
        tf.margin_top = Emu(80000)
        tf.margin_right = Emu(120000)
        tf.margin_bottom = Emu(80000)

        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(1)
            p.space_before = Pt(0)
            add_run(p, line, font_name='Consolas', font_size=Pt(10),
                    color=RGBColor(0xCC, 0xDD, 0xEE))

        y_pos += code_height + Emu(150000)


def rebuild_content_slide(slide, title_text, body_content, title_size=Pt(36)):
    """Rebuild a content slide with proper sizing."""
    clear_slide_shapes(slide, keep_placeholders=True)

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title_text
            for run in shape.text_frame.paragraphs[0].runs:
                run.font.size = title_size
                run.font.bold = True
            break

    txb = add_textbox(slide, Emu(500000), Emu(1100000), Emu(11200000), Emu(5500000))
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for item in body_content:
        if isinstance(item, str):
            text, bold, size, color, indent = item, False, Pt(18), None, 0
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            size = item[2] if len(item) > 2 else Pt(18)
            color = item[3] if len(item) > 3 else None
            indent = item[4] if len(item) > 4 else 0

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        if indent > 0:
            p.level = indent

        p.space_after = Pt(2)
        p.space_before = Pt(1)
        add_run(p, text, font_size=size, bold=bold, color=color)

    return slide


# ═══════════════════════════════════════
# MAIN
# ═══════════════════════════════════════

def main():
    print(f"Loading {SRC}...")
    prs = Presentation(SRC)
    print(f"  {len(prs.slides)} slides")

    # ═══════════════════════════════════
    # FIX 1: Trainer slide (slide 6) overflow
    # ═══════════════════════════════════
    print("\n--- Fix 1: Trainer slide positioning ---")
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        # Fix right overflow on subtitle textbox
        r = shape.left + shape.width
        if r > SLIDE_W:
            new_width = SLIDE_W - shape.left - Emu(100000)
            if new_width > 0:
                shape.width = new_width
                print(f"  Fixed width of '{shape.name}': trimmed to fit")

        # Also shrink the right-column textboxes if they extend too far
        if shape.left + shape.width > SLIDE_W:
            shape.width = SLIDE_W - shape.left - Emu(50000)
            print(f"  Fixed '{shape.name}' right edge")

    print("  Trainer slide fixed")

    # ═══════════════════════════════════
    # FIX 2: Code slides overflow (100, 102, 103)
    # ═══════════════════════════════════
    print("\n--- Fix 2: Code slides overflow ---")

    # Find slides by title
    slide_map = {}
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    slide_map[t[:60]] = i
                    break

    # Slide: "Building an LLM-SAST Scanner (Python + Ollama)"
    idx = slide_map.get("Building an LLM-SAST Scanner (Python + Ollama)")
    if idx is not None:
        print(f"  Rebuilding slide {idx+1}: LLM-SAST Scanner...")
        rebuild_code_slide(prs.slides[idx],
            "Building an LLM-SAST Scanner (Python + Ollama)",
            [("llm_sast_scanner.py:", """import ollama, json, sys, pathlib

SYSTEM_PROMPT = pathlib.Path("prompt.txt").read_text()

def scan_file(filepath):
    code = pathlib.Path(filepath).read_text()
    response = ollama.chat(
        model='qwen2.5-coder:7b',
        messages=[
            {'role': 'system', 'content': SYSTEM_PROMPT},
            {'role': 'user',
             'content': f'Analyze:\\n```\\n{code}\\n```'}
        ], format='json')
    return json.loads(response['message']['content'])

if __name__ == '__main__':
    results = scan_file(sys.argv[1])
    for v in results.get('vulnerabilities', []):
        print(f"[{v['severity']}] {v['type']} @ line {v['line']}")
        print(f"  -> {v.get('explanation','')[:100]}")""")])

    # Slide: "Using Claude API for Cloud-Based LLM-SAST"
    idx = slide_map.get("Using Claude API for Cloud-Based LLM-SAST")
    if idx is not None:
        print(f"  Rebuilding slide {idx+1}: Claude API SAST...")
        rebuild_code_slide(prs.slides[idx],
            "Using Claude API for Cloud-Based LLM-SAST",
            [("claude_sast.py:", """import anthropic, json, sys, pathlib

client = anthropic.Anthropic()  # ANTHROPIC_API_KEY env var
PROMPT = pathlib.Path("prompt.txt").read_text()

def scan_file(filepath):
    code = pathlib.Path(filepath).read_text()
    msg = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4096, system=PROMPT,
        messages=[{'role':'user',
            'content':f'Analyze:\\n```\\n{code}\\n```'}])
    return json.loads(msg.content[0].text)

# pip install anthropic
# export ANTHROPIC_API_KEY="sk-ant-..."
# python claude_sast.py target.php""")])

    # Slide: "CI/CD Integration: GitHub Actions"  — SPLIT into 2 slides
    idx = slide_map.get("CI/CD Integration: GitHub Actions")
    if idx is not None:
        print(f"  Rebuilding slide {idx+1}: CI/CD Integration...")
        rebuild_code_slide(prs.slides[idx],
            "CI/CD Integration: GitHub Actions",
            [(".github/workflows/llm-sast.yml:", """name: LLM Security Scan
on: [pull_request]
jobs:
  llm-sast:
    runs-on: ubuntu-latest
    steps:
      - uses: actions/checkout@v4
      - uses: actions/setup-python@v5
        with: { python-version: '3.12' }
      - run: pip install anthropic
      - name: Scan changed files
        env:
          ANTHROPIC_API_KEY: ${{ secrets.ANTHROPIC_API_KEY }}
        run: |
          git diff --name-only origin/main -- '*.php' '*.py' |
          while read f; do
            python llm_sast_scanner.py "$f" >> report.md
          done
      - name: Comment on PR
        uses: marocchino/sticky-pull-request-comment@v2
        with:
          path: report.md""")])

    # ═══════════════════════════════════
    # FIX 3: Prompt Engineering slide — check and fix
    # ═══════════════════════════════════
    idx = slide_map.get("Prompt Engineering: Security Analysis Prompt")
    if idx is not None:
        print(f"  Rebuilding slide {idx+1}: Prompt Engineering...")
        rebuild_code_slide(prs.slides[idx],
            "Prompt Engineering: Security Analysis Prompt",
            [("System Prompt (save as prompt.txt):", """You are an expert application security engineer.
Analyze the following source code for vulnerabilities.

For each finding, provide JSON with:
- type: vulnerability type (SQL Injection, XSS, etc)
- cwe: CWE ID number
- severity: Critical / High / Medium / Low
- line: affected line number(s)
- explanation: attack vector description
- fix: secure code replacement

Focus on OWASP Top 10 2025. Report only real,
exploitable vulnerabilities. No style or quality issues.
Respond ONLY with valid JSON: {"vulnerabilities":[...]}""")])

    # ═══════════════════════════════════
    # FIX 4: Running the scanner — fix
    # ═══════════════════════════════════
    idx = slide_map.get("Running the LLM-SAST Scanner")
    if idx is not None:
        print(f"  Rebuilding slide {idx+1}: Running scanner...")
        rebuild_code_slide(prs.slides[idx],
            "Running the LLM-SAST Scanner",
            [("Usage & Example Output:", """# Scan a single file
$ python llm_sast_scanner.py vulnerable_login.php

[CRITICAL] SQL Injection @ line 23
  -> $_GET['id'] concatenated into query without
     parameterization. Use PDO prepared statements.

[HIGH] Stored XSS @ line 45
  -> User data echoed without htmlspecialchars().
     Apply output encoding before HTML rendering.

[MEDIUM] IDOR @ line 67
  -> user_id from request used without ownership
     check. Validate resource belongs to session user.

# Scan a whole project
$ find src/ -name '*.php' | xargs -I{} python \\
    llm_sast_scanner.py {}""")])

    # ═══════════════════════════════════
    # FIX 5: Comparison table — rebuild with better layout
    # ═══════════════════════════════════
    idx = slide_map.get("Traditional SAST vs LLM-Powered SAST")
    if idx is not None:
        print(f"  Rebuilding slide {idx+1}: Comparison table...")
        rebuild_content_slide(prs.slides[idx],
            "Traditional SAST vs LLM-Powered SAST", [
                ("", False, Pt(6)),
                ("Detection", True, Pt(17), XECURA_DARK_BLUE),
                ("  Traditional: Regex / AST pattern matching", False, Pt(15)),
                ("  LLM: Semantic code understanding & reasoning", False, Pt(15), XECURA_MEDIUM_BLUE),
                ("False Positives", True, Pt(17), XECURA_DARK_BLUE),
                ("  Traditional: High (30-70%)", False, Pt(15)),
                ("  LLM: Low — filters with context awareness", False, Pt(15), XECURA_MEDIUM_BLUE),
                ("Business Logic", True, Pt(17), XECURA_DARK_BLUE),
                ("  Traditional: Cannot detect logic flaws", False, Pt(15)),
                ("  LLM: Can reason about design & logic issues", False, Pt(15), XECURA_MEDIUM_BLUE),
                ("Fix Suggestions", True, Pt(17), XECURA_DARK_BLUE),
                ("  Traditional: Generic CWE templates", False, Pt(15)),
                ("  LLM: Context-aware code patches", False, Pt(15), XECURA_MEDIUM_BLUE),
                ("Speed", True, Pt(17), XECURA_DARK_BLUE),
                ("  Traditional: Very fast (milliseconds)", False, Pt(15)),
                ("  LLM: Slower (seconds per file)", False, Pt(15), XECURA_MEDIUM_BLUE),
                ("Offline", True, Pt(17), XECURA_DARK_BLUE),
                ("  Traditional: Full offline support", False, Pt(15)),
                ("  LLM: Ollama = offline, Cloud APIs = online", False, Pt(15), XECURA_MEDIUM_BLUE),
                ("", False, Pt(6)),
                ("Best Practice: Use BOTH — traditional for speed, LLM for depth", True, Pt(16), XECURA_RED),
            ], title_size=Pt(32))

    # ═══════════════════════════════════
    # FIX 6: Update "Tools & Models" to reference qwen2.5-coder:7b
    # ═══════════════════════════════════
    idx = slide_map.get("LLM-Powered SAST Tools & Models")
    if idx is not None:
        print(f"  Rebuilding slide {idx+1}: Tools & Models...")
        rebuild_content_slide(prs.slides[idx],
            "LLM-Powered SAST Tools & Models", [
                ("Commercial / SaaS:", True, Pt(20), XECURA_DARK_BLUE),
                ("GitHub Copilot Autofix — integrated in GitHub Advanced Security", False, Pt(16), None, 1),
                ("Snyk DeepCode AI — AI code analysis in Snyk platform", False, Pt(16), None, 1),
                ("Amazon CodeGuru — ML-based code review for Java/Python", False, Pt(16), None, 1),
                ("Semgrep + AI Assistant — rule-based + LLM-assisted", False, Pt(16), None, 1),
                ("", False, Pt(8)),
                ("Open Source / Self-Hosted:", True, Pt(20), XECURA_DARK_BLUE),
                ("Claude API — Anthropic, strong at code analysis", False, Pt(16), None, 1),
                ("GPT-4 / OpenAI API — effective with good prompts", False, Pt(16), None, 1),
                ("Ollama + Qwen2.5-Coder — fully offline, self-hosted", False, Pt(16), None, 1),
                ("Ollama + DeepSeek-Coder / CodeLlama — alternatives", False, Pt(16), None, 1),
                ("", False, Pt(8)),
                ("For this training we use:", True, Pt(18), XECURA_RED),
                ("  Ollama + qwen2.5-coder:7b (offline, runs on this machine)", False, Pt(16), XECURA_RED),
            ])

    # ═══════════════════════════════════
    # FIX 7: Setup slide — update to qwen2.5-coder:7b
    # ═══════════════════════════════════
    idx = slide_map.get("Setup: Ollama for Offline LLM-SAST")
    if idx is not None:
        print(f"  Rebuilding slide {idx+1}: Ollama Setup...")
        rebuild_content_slide(prs.slides[idx],
            "Setup: Ollama for Offline LLM-SAST", [
                ("Step 1: Install Ollama", True, Pt(20), XECURA_DARK_BLUE),
                ("  curl -fsSL https://ollama.com/install.sh | sh", False, Pt(14)),
                ("  # macOS/Windows: download from https://ollama.com", False, Pt(14)),
                ("", False, Pt(6)),
                ("Step 2: Pull a Code Analysis Model", True, Pt(20), XECURA_DARK_BLUE),
                ("  ollama pull qwen2.5-coder:7b", False, Pt(14)),
                ("  # Alternatives:", False, Pt(14)),
                ("  # ollama pull deepseek-coder-v2:16b  (needs 16GB+ RAM)", False, Pt(14)),
                ("  # ollama pull codellama:13b", False, Pt(14)),
                ("", False, Pt(6)),
                ("Step 3: Verify", True, Pt(20), XECURA_DARK_BLUE),
                ("  ollama list", False, Pt(14)),
                ("  ollama run qwen2.5-coder:7b 'Hello, test'", False, Pt(14)),
                ("", False, Pt(6)),
                ("Requirements:", True, Pt(18), XECURA_RED),
                ("  qwen2.5-coder:7b needs ~5GB RAM — runs on most laptops", False, Pt(14)),
                ("  GPU recommended but not required", False, Pt(14)),
            ])

    # ═══════════════════════════════════
    # FIX 8: Rearrange — move LLM hands-on into Module 8
    # ═══════════════════════════════════
    print("\n--- Fix 8: Rearrange hands-on slides ---")

    # First, find the "Hands-On: LLM vs Traditional SAST" slide
    llm_handson_idx = None
    module8_divider_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if 'Hands-On: LLM vs Traditional SAST' in t:
                    llm_handson_idx = i
                if 'Module 8:' in t and 'Hands-On Lab' in t:
                    module8_divider_idx = i

    if llm_handson_idx is not None and module8_divider_idx is not None:
        print(f"  LLM Hands-On at slide {llm_handson_idx+1}, Module 8 divider at {module8_divider_idx+1}")
        # Move LLM hands-on to right after Module 8 divider
        # It should go after the existing hands-on slides (DVWA, Kali)
        # Module 8 divider is at module8_divider_idx, followed by existing hands-on slides
        # We want LLM hands-on AFTER those
        target = module8_divider_idx + 3  # after divider + kali + dvwa
        if llm_handson_idx < target:
            target -= 1  # adjust if source is before target
        move_slide(prs, llm_handson_idx, target)
        print(f"  Moved LLM Hands-On to position {target+1}")
    else:
        print(f"  WARNING: Could not find slides (llm={llm_handson_idx}, m8={module8_divider_idx})")

    # ═══════════════════════════════════
    # FIX 9: Update OWASP 2021 → 2025 throughout
    # ═══════════════════════════════════
    print("\n--- Fix 9: Update OWASP 2021 → 2025 ---")

    # OWASP 2025 mapping (what changed from 2021)
    owasp_renames = {
        "A01:2021-Broken Access Control": "A01:2025-Broken Access Control",
        "A02:2021-Cryptographic Failures": "A04:2025-Cryptographic Failures",
        "A03:2021-Injection": "A05:2025-Injection",
        "A04:2021-Insecure Design": "A06:2025-Insecure Design",
        "A05:2021-Security Misconfiguration": "A02:2025-Security Misconfiguration",
        "A06:2021-Vulnerable and Outdated Components": "A03:2025-Software Supply Chain Failures",
        "A07:2021-Identification & Authentication Failures": "A07:2025-Authentication Failures",
        "A08:2021-Software and Data Integrity Failures": "A08:2025-Software or Data Integrity Failures",
        "A09:2021-Security Logging & Monitoring Failures": "A09:2025-Security Logging and Alerting Failures",
        "A10:2021-Server-Side Request Forgery (SSRF)": "A10:2025-Mishandling of Exceptional Conditions",
    }

    # Update the Module 3 divider title
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "OWASP Top 10 (2021)" in run.text:
                            run.text = run.text.replace("OWASP Top 10 (2021)", "OWASP Top 10 (2025)")
                            print(f"  Slide {i+1}: Updated module title to 2025")
                        if "2021" in run.text and "OWASP" in para.text:
                            run.text = run.text.replace("2021", "2025")

    # Update individual OWASP slide titles
    rename_count = 0
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                full = shape.text_frame.text
                for old, new in owasp_renames.items():
                    if old in full:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if old in run.text:
                                    run.text = run.text.replace(old, new)
                                    rename_count += 1

    print(f"  Updated {rename_count} OWASP title references")

    # Update the OWASP intro/description slides
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "OWASP Top 10 (2021)" in run.text:
                            run.text = run.text.replace("OWASP Top 10 (2021)", "OWASP Top 10 (2025)")
                        if "OWASP Top 10 2021" in run.text:
                            run.text = run.text.replace("OWASP Top 10 2021", "OWASP Top 10 2025")

    # Update outline slide
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if "OWASP Top 10 (2021)" in para.text:
                        for run in para.runs:
                            run.text = run.text.replace("2021", "2025")
                            print(f"  Slide {i+1}: Updated outline OWASP ref")

    # ═══════════════════════════════════
    # SAVE
    # ═══════════════════════════════════
    print(f"\nSaving to {DST}...")
    prs.save(DST)

    # Verify final structure
    prs2 = Presentation(DST)
    print(f"\n=== FINAL SLIDE MAP ({len(prs2.slides)} slides) ===")
    for i, slide in enumerate(prs2.slides):
        title = ''
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t and len(t) > 3 and t != str(i+1):
                    title = t[:80]
                    break
        if any(k in title for k in ['Module ', 'INTROS', 'Outline', 'WHO AM', 'QUIZ', 'Any question',
                                      'Survey', 'LLM', 'Hands-On', 'OWASP', 'A0']):
            print(f'  Slide {i+1:3d}: {title}')

    print("\nDone!")


if __name__ == '__main__':
    main()
