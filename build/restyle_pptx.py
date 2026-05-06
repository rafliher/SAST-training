#!/usr/bin/env python3
"""
Restyle the Secure Coding PPTX to match Xecura branding.
Two-phase approach:
  Phase 1: python-pptx for colors, text, and shape-level edits
  Phase 2: ZIP-level surgery to replace image files and fix master XML
"""

from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from lxml import etree
import zipfile
import shutil
import os
import re
import io
import tempfile

# ── Xecura Color Palette ──
XECURA_DARK_BLUE  = RGBColor(0x1D, 0x5C, 0xA8)  # #1D5CA8
XECURA_MEDIUM_BLUE = RGBColor(0x0F, 0x81, 0xC9)  # #0F81C9
XECURA_CYAN       = RGBColor(0x00, 0xBA, 0xFF)    # #00BAFF
XECURA_RED        = RGBColor(0xEF, 0x35, 0x4F)    # #EF354F

COLOR_MAP = {
    '002060': '1D5CA8',
    'FF0066': 'EF354F',
    '00B0F0': '00BAFF',
    '0070C0': '0F81C9',
    '0563C1': '0F81C9',
}

# ── Paths ──
SRC = 'day 1/Secure Coding in Software Development [19-21 Agustus]_v1.pptx'
DST = 'day 1/Secure Coding - Xecura Branded.pptx'
XECURA_FULL_LOGO = 'assets/xecura_bg.png'    # Full logo with "XECURA" text
XECURA_ICON      = 'assets/xecura_logo.png'  # Shield icon only


def hex_color(rgb):
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


# ═══════════════════════════════════════════
# PHASE 1: python-pptx based edits
# ═══════════════════════════════════════════

def remap_colors_in_xml(element):
    """Find and remap all srgbClr values in an XML element tree."""
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    for srgb in element.iter(f'{ns}srgbClr'):
        val = srgb.get('val', '').upper()
        if val in COLOR_MAP:
            srgb.set('val', COLOR_MAP[val])


def update_branding_text(shape):
    """Replace Edutech references with Xecura."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        full_text = para.text
        if 'Edutech' in full_text or 'edutech' in full_text.lower():
            for run in para.runs:
                run.text = run.text.replace('Edutech', 'Xecura')
                run.text = run.text.replace('edutech', 'xecura')
                run.text = run.text.replace('© 2017, ', '© 2025, ')
                run.text = run.text.replace('Solution program', 'Cyber Security')
        if 'edutechsolution' in full_text.lower():
            for run in para.runs:
                run.text = run.text.replace('edutechsolution.co.id', 'xecura.id')
                run.text = run.text.replace('https://xecura.id/mV', 'https://xecura.id')


def phase1_pptx_edits():
    """Use python-pptx to fix colors and text across all slides, masters, layouts."""
    print("Phase 1: python-pptx color/text edits...")
    prs = Presentation(SRC)

    # Fix slide masters + layouts
    for master in prs.slide_masters:
        remap_colors_in_xml(master._element)
        # Fix branding text in master shapes
        for shape in master.shapes:
            remap_colors_in_xml(shape._element)
            update_branding_text(shape)
            # Direct fix for "EDUTECH" sidebar text in master
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip() == 'EDUTECH':
                            run.text = 'XECURA'
                            print(f"  Fixed master text: EDUTECH → XECURA")
                        if run.text.strip() == 'MODUL MATERI':
                            run.text = 'TRAINING MODULE'
                            print(f"  Fixed master text: MODUL MATERI → TRAINING MODULE")
        for layout in master.slide_layouts:
            remap_colors_in_xml(layout._element)
            for shape in layout.shapes:
                remap_colors_in_xml(shape._element)
                update_branding_text(shape)
            for shape in layout.placeholders:
                remap_colors_in_xml(shape._element)

    # Fix every slide
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_num % 20 == 1:
            print(f"  Slides {slide_num}-{min(slide_num+19, total)}/{total}...")

        # Colors in all shapes via XML
        for shape in slide.shapes:
            remap_colors_in_xml(shape._element)
            update_branding_text(shape)

    # Save intermediate
    prs.save(DST)
    print(f"  Saved intermediate to {DST}")


# ═══════════════════════════════════════════
# PHASE 2: ZIP-level image replacement
# ═══════════════════════════════════════════

def phase2_zip_surgery():
    """Directly replace image files and fix XML inside the PPTX zip."""
    print("Phase 2: ZIP-level image replacement...")

    with open(XECURA_FULL_LOGO, 'rb') as f:
        xecura_full_logo_bytes = f.read()
    with open(XECURA_ICON, 'rb') as f:
        xecura_icon_bytes = f.read()

    # Images to replace:
    # image1.png = Edutech logo in slide master (appears on EVERY slide)
    # image3.png = Edutech logo on title slides 1 & 60
    image_replacements = {
        'ppt/media/image1.png': xecura_full_logo_bytes,   # Master logo → Xecura full
        'ppt/media/image2.png': xecura_icon_bytes,         # Title slide Picture 5 → Xecura icon
        'ppt/media/image3.png': xecura_full_logo_bytes,   # Title slide logo → Xecura full
    }

    # XML text replacements (for master/layout XML that python-pptx might miss)
    xml_text_fixes = {
        'Edutech': 'Xecura',
        'edutech': 'xecura',
    }

    # Color replacements in ALL XML files (catch anything python-pptx missed)
    color_fixes = {
        'val="002060"': 'val="1D5CA8"',
        'val="FF0066"': 'val="EF354F"',
        'val="00B0F0"': 'val="00BAFF"',
        'val="0070C0"': 'val="0F81C9"',
        'val="0563C1"': 'val="0F81C9"',
    }

    tmp_path = DST + '.tmp'
    replaced_images = set()
    fixed_xml = set()

    with zipfile.ZipFile(DST, 'r') as zin:
        with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)

                # Replace images
                if item.filename in image_replacements:
                    data = image_replacements[item.filename]
                    replaced_images.add(item.filename)
                    print(f"  Replaced image: {item.filename}")

                # Fix XML files
                elif item.filename.endswith('.xml') or item.filename.endswith('.rels'):
                    text = data.decode('utf-8')
                    changed = False

                    # Fix colors
                    for old, new in color_fixes.items():
                        if old in text:
                            count = text.count(old)
                            text = text.replace(old, new)
                            changed = True
                            print(f"  Fixed {count}x color {old} in {item.filename}")

                    # Fix branding text in XML
                    for old, new in xml_text_fixes.items():
                        if old in text:
                            text = text.replace(old, new)
                            changed = True

                    if changed:
                        data = text.encode('utf-8')
                        fixed_xml.add(item.filename)

                zout.writestr(item, data)

    # Replace original with fixed version
    shutil.move(tmp_path, DST)
    print(f"  Replaced {len(replaced_images)} images, fixed {len(fixed_xml)} XML files")


# ═══════════════════════════════════════════
# PHASE 3: Verification
# ═══════════════════════════════════════════

def phase3_verify():
    """Verify no Edutech remnants remain."""
    print("Phase 3: Verification...")
    issues = []

    with zipfile.ZipFile(DST, 'r') as z:
        for name in z.namelist():
            # Check images
            if name in ('ppt/media/image1.png', 'ppt/media/image3.png'):
                data = z.read(name)
                with open(XECURA_FULL_LOGO, 'rb') as f:
                    expected = f.read()
                if data == expected:
                    print(f"  OK: {name} = Xecura logo")
                else:
                    issues.append(f"  FAIL: {name} was NOT replaced!")

            # Check XML for remnants
            if name.endswith('.xml'):
                content = z.read(name).decode('utf-8')

                # Check for Edutech text
                if 'dutech' in content.lower():
                    matches = re.findall(r'[^<>]*[Ee]dutech[^<>]*', content)
                    for m in matches[:2]:
                        issues.append(f"  REMAINING TEXT in {name}: {m.strip()[:80]}")

                # Check for old colors
                for old_color in ['002060', 'FF0066', '00B0F0', '0070C0']:
                    count = content.count(f'val="{old_color}"')
                    if count > 0:
                        issues.append(f"  REMAINING COLOR {old_color} x{count} in {name}")

    if issues:
        print("  ISSUES FOUND:")
        for issue in issues:
            print(f"    {issue}")
    else:
        print("  ALL CLEAN - no Edutech remnants found!")

    return len(issues) == 0


def main():
    phase1_pptx_edits()
    phase2_zip_surgery()
    ok = phase3_verify()
    print()
    if ok:
        print(f"SUCCESS: Restyled PPTX saved to: {DST}")
    else:
        print(f"WARNING: Some issues remain. Check output above.")
    print(f"File size: {os.path.getsize(DST) / 1024 / 1024:.1f} MB")


if __name__ == '__main__':
    main()
